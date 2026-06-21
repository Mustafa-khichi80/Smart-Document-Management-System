"""
PPTX Converter
Converts PowerPoint files to PDF and images.
Uses PyMuPDF for PDF to image conversion (no Poppler needed).
"""
import os
import asyncio
import subprocess
import shutil


async def convert_pptx(input_path: str, output_dir: str, target_format: str) -> dict:
    """PPTX converter supporting PDF and image outputs."""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pptx, input_path, output_dir, target_format)


def _process_pptx(input_path: str, output_dir: str, target_format: str) -> dict:
    try:
        filename = os.path.basename(input_path)
        name, ext = os.path.splitext(filename)
        output_filename = f"{name}.{target_format}"
        output_path = os.path.join(output_dir, output_filename)

        if target_format == 'pdf':
            return _pptx_to_pdf(input_path, output_path, output_filename)
        elif target_format in ['png', 'jpg', 'jpeg']:
            return _pptx_to_images(input_path, output_dir, target_format, name)
        elif target_format == 'txt':
            return _pptx_to_txt(input_path, output_path, output_filename)
        else:
            return {"success": False, "error": f"Unsupported target format for PPTX: {target_format}"}

    except Exception as e:
        return {"success": False, "error": f"PPTX conversion error: {str(e)}"}


def _pptx_to_pdf(input_path: str, output_path: str, output_filename: str) -> dict:
    """Convert PPTX to PDF using LibreOffice, PowerPoint COM, or ReportLab fallback."""
    
    # Method 1: Try LibreOffice (cross-platform)
    try:
        libreoffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "soffice",
            "/usr/bin/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        ]
        
        soffice_path = None
        for path in libreoffice_paths:
            if os.path.exists(path) or shutil.which(path):
                soffice_path = path
                break
        
        if soffice_path:
            out_dir = os.path.dirname(output_path)
            result = subprocess.run([
                soffice_path,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", out_dir,
                input_path
            ], capture_output=True, text=True, timeout=180)
            
            expected_output = os.path.join(out_dir, os.path.splitext(os.path.basename(input_path))[0] + ".pdf")
            if os.path.exists(expected_output):
                if expected_output != output_path:
                    shutil.move(expected_output, output_path)
                return {"success": True, "output_path": output_path, "filename": output_filename}
    except Exception as e:
        print(f"[PPTX->PDF] LibreOffice failed: {e}")
    
    # Method 2: Try PowerPoint COM (Windows only)
    try:
        import platform
        if platform.system() == 'Windows':
            import comtypes.client
            
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            
            presentation = powerpoint.Presentations.Open(os.path.abspath(input_path))
            presentation.SaveAs(os.path.abspath(output_path), 32)  # 32 = ppSaveAsPDF
            presentation.Close()
            powerpoint.Quit()
            
            if os.path.exists(output_path):
                return {"success": True, "output_path": output_path, "filename": output_filename}
    except Exception as e:
        print(f"[PPTX->PDF] PowerPoint COM failed: {e}")
    
    # Method 3: Python-only fallback (ReportLab) - used when PowerPoint/LibreOffice are missing (e.g. on Vercel)
    try:
        from pptx import Presentation
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        
        prs = Presentation(input_path)
        pdf = SimpleDocTemplate(output_path, pagesize=landscape(letter),
                                rightMargin=40, leftMargin=40,
                                topMargin=40, bottomMargin=40)
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading2'],
            fontName='Helvetica-Bold',
            fontSize=18,
            leading=22,
            spaceAfter=15
        )
        body_style = ParagraphStyle(
            'SlideBody',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=12,
            leading=16,
            spaceAfter=10
        )
        
        story = []
        for idx, slide in enumerate(prs.slides):
            if idx > 0:
                story.append(PageBreak())
                
            story.append(Paragraph(f"Slide {idx+1}", title_style))
            story.append(Spacer(1, 10))
            
            p_counter = 0
            # Extract text from shapes with styling preservation
            for shape in slide.shapes:
                # 1. Text Frame Shapes
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        p_html = ""
                        for run in paragraph.runs:
                            text = run.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            if not text:
                                continue
                            
                            tags_start = ""
                            tags_end = ""
                            
                            font_attrs = []
                            if run.font.size:
                                font_attrs.append(f'size="{run.font.size.pt}"')
                            if run.font.color and run.font.color.type == 1:
                                hex_color = f"#{run.font.color.rgb}"
                                font_attrs.append(f'color="{hex_color}"')
                            if run.font.name:
                                font_name = run.font.name.lower()
                                if "times" in font_name:
                                    font_attrs.append('name="Times-Roman"')
                                elif "courier" in font_name:
                                    font_attrs.append('name="Courier"')
                                else:
                                    font_attrs.append('name="Helvetica"')
                                    
                            if font_attrs:
                                tags_start += f'<font {" ".join(font_attrs)}>'
                                tags_end = '</font>' + tags_end
                                
                            if run.font.bold:
                                tags_start += '<b>'
                                tags_end = '</b>' + tags_end
                            if run.font.italic:
                                tags_start += '<i>'
                                tags_end = '</i>' + tags_end
                            if run.font.underline:
                                tags_start += '<u>'
                                tags_end = '</u>' + tags_end
                                
                            p_html += f"{tags_start}{text}{tags_end}"
                        
                        if p_html.strip():
                            align_val = 0
                            if paragraph.alignment == 2:
                                align_val = 1
                            elif paragraph.alignment == 3:
                                align_val = 2
                            elif paragraph.alignment == 4:
                                align_val = 4
                                
                            p_style = ParagraphStyle(
                                f'SStyle_{idx}_{p_counter}',
                                parent=body_style,
                                alignment=align_val
                            )
                            story.append(Paragraph(p_html, p_style))
                            p_counter += 1
                
                # 2. Table Shapes
                elif hasattr(shape, "has_table") and shape.has_table:
                    from reportlab.platypus import Table, TableStyle
                    from reportlab.lib import colors
                    
                    table_data = []
                    t_styles = [
                        ('VALIGN', (0,0), (-1,-1), 'TOP'),
                        ('GRID', (0,0), (-1,-1), 0.5, colors.grey),
                        ('TOPPADDING', (0,0), (-1,-1), 4),
                        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
                        ('LEFTPADDING', (0,0), (-1,-1), 4),
                        ('RIGHTPADDING', (0,0), (-1,-1), 4),
                    ]
                    
                    t_obj = shape.table
                    
                    for r_idx, row in enumerate(t_obj.rows):
                        row_data = []
                        for c_idx, cell in enumerate(row.cells):
                            cell_paragraphs = []
                            for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                                p_html = ""
                                for run in paragraph.runs:
                                    text = run.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                    if not text:
                                        continue
                                    
                                    tags_start = ""
                                    tags_end = ""
                                    
                                    font_attrs = []
                                    if run.font.size:
                                        font_attrs.append(f'size="{run.font.size.pt}"')
                                    if run.font.color and run.font.color.type == 1:
                                        try:
                                            hex_color = f"#{run.font.color.rgb}"
                                            font_attrs.append(f'color="{hex_color}"')
                                        except:
                                            pass
                                    if run.font.name:
                                        font_name = run.font.name.lower()
                                        if "times" in font_name:
                                            font_attrs.append('name="Times-Roman"')
                                        elif "courier" in font_name:
                                            font_attrs.append('name="Courier"')
                                        else:
                                            font_attrs.append('name="Helvetica"')
                                            
                                    if font_attrs:
                                        tags_start += f'<font {" ".join(font_attrs)}>'
                                        tags_end = '</font>' + tags_end
                                        
                                    if run.font.bold:
                                        tags_start += '<b>'
                                        tags_end = '</b>' + tags_end
                                    if run.font.italic:
                                        tags_start += '<i>'
                                        tags_end = '</i>' + tags_end
                                    if run.font.underline:
                                        tags_start += '<u>'
                                        tags_end = '</u>' + tags_end
                                        
                                    p_html += f"{tags_start}{text}{tags_end}"
                                
                                align_val = 0
                                if paragraph.alignment == 2:
                                    align_val = 1
                                elif paragraph.alignment == 3:
                                    align_val = 2
                                elif paragraph.alignment == 4:
                                    align_val = 4
                                    
                                p_style = ParagraphStyle(
                                    f'TStyle_{idx}_{r_idx}_{c_idx}_{p_idx}',
                                    parent=body_style,
                                    alignment=align_val
                                )
                                cell_paragraphs.append(Paragraph(p_html or "&nbsp;", p_style))
                            row_data.append(cell_paragraphs)
                        table_data.append(row_data)
                    
                    if table_data:
                        rl_table = Table(table_data)
                        rl_table.setStyle(TableStyle(t_styles))
                        story.append(Spacer(1, 10))
                        story.append(rl_table)
                        story.append(Spacer(1, 10))
            
        pdf.build(story)
        if os.path.exists(output_path):
            return {
                "success": True, 
                "output_path": output_path, 
                "filename": output_filename,
                "note": "Converted using ReportLab fallback (PowerPoint/LibreOffice missing)"
            }
    except Exception as fallback_err:
        print(f"[PPTX->PDF] Fallback failed: {fallback_err}")
    
    return {
        "success": False, 
        "error": "PDF conversion requires PowerPoint or LibreOffice. Please install one of them."
    }


def _pptx_to_images(input_path: str, output_dir: str, target_format: str, name: str) -> dict:
    """Convert PPTX to images using PyMuPDF (no Poppler needed)."""
    import zipfile
    
    # First convert to PDF
    pdf_path = os.path.join(output_dir, f"{name}_temp.pdf")
    pdf_result = _pptx_to_pdf(input_path, pdf_path, f"{name}.pdf")
    
    if not pdf_result["success"]:
        # Fallback: Extract embedded images from PPTX directly
        try:
            images_extracted = []
            with zipfile.ZipFile(input_path, 'r') as pptx:
                for i, item in enumerate(pptx.namelist()):
                    if item.startswith('ppt/media/') and any(item.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif']):
                        img_data = pptx.read(item)
                        original_ext = os.path.splitext(item)[1]
                        img_filename = f"{name}_image_{i+1}{original_ext}"
                        img_path = os.path.join(output_dir, img_filename)
                        with open(img_path, 'wb') as f:
                            f.write(img_data)
                        images_extracted.append(img_filename)
            
            if images_extracted:
                return {
                    "success": True, 
                    "output_path": os.path.join(output_dir, images_extracted[0]), 
                    "filename": images_extracted[0],
                    "note": f"Extracted {len(images_extracted)} images from presentation"
                }
            return {"success": False, "error": "No images found in presentation and PDF conversion failed"}
        except Exception as e:
            return {"success": False, "error": f"Image extraction error: {str(e)}"}
    
    # Convert PDF to images using PyMuPDF (fitz) - NO POPPLER NEEDED!
    try:
        import fitz  # PyMuPDF
        
        pdf_doc = fitz.open(pdf_path)
        output_files = []
        
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            # Higher resolution for better quality
            mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for 144 DPI
            pix = page.get_pixmap(matrix=mat)
            
            img_filename = f"{name}_slide_{page_num + 1}.{target_format}"
            img_path = os.path.join(output_dir, img_filename)
            
            if target_format.lower() in ['jpg', 'jpeg']:
                pix.save(img_path, output="jpeg", jpg_quality=95)
            else:
                pix.save(img_path)
            
            output_files.append(img_filename)
        
        pdf_doc.close()
        
        # Clean up temp PDF
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        
        if output_files:
            return {
                "success": True, 
                "output_path": os.path.join(output_dir, output_files[0]), 
                "filename": output_files[0],
                "all_files": output_files,
                "note": f"Created {len(output_files)} slide images"
            }
        
        return {"success": False, "error": "Could not convert slide"}
        
    except ImportError:
        return {"success": False, "error": "PyMuPDF not installed. Run 'pip install PyMuPDF'"}
    except Exception as e:
        # Clean up on error
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        return {"success": False, "error": f"PDF to image conversion error: {str(e)}"}


def _pptx_to_txt(input_path: str, output_path: str, output_filename: str) -> dict:
    """Extract text from PPTX."""
    try:
        from pptx import Presentation
    except ImportError:
        return {"success": False, "error": "python-pptx not installed. Run 'pip install python-pptx'"}
    
    prs = Presentation(input_path)
    text_content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        text_content.append(f"=== Slide {slide_num} ===\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text)
        text_content.append("\n")
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write("\n".join(text_content))
    
    return {"success": True, "output_path": output_path, "filename": output_filename}
