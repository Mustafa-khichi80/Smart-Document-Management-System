"""
Utility functions for Smart Document Management.
"""
import shutil
import os
import re


def get_output_dir():
    """Get or create the output directory for converted files."""
    IS_VERCEL = os.environ.get("VERCEL") == "1"
    if IS_VERCEL:
        output_dir = "/tmp/converted_files"
    else:
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        output_dir = os.path.join(base_dir, "converted_files")
        
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    return output_dir

def clean_filename(filename: str) -> str:
    """
    Sanitize filename for security.
    Removes path injection, special characters and Turkish chars.
    """
    if not filename:
        return "unnamed_file"
    
    filename = filename.replace('\\', '/')
    filename = os.path.basename(filename)
    
    filename = filename.replace('\x00', '')
    filename = re.sub(r'[\x00-\x1f\x7f]', '', filename)
    
    dangerous_chars = '<>:"|?*'
    for char in dangerous_chars:
        filename = filename.replace(char, '')
    
    tr_replacements = {
        'ı': 'i', 'ğ': 'g', 'ü': 'u', 'ş': 's', 'ö': 'o', 'ç': 'c',
        'İ': 'I', 'Ğ': 'G', 'Ü': 'U', 'Ş': 'S', 'Ö': 'O', 'Ç': 'C'
    }
    for char, replacement in tr_replacements.items():
        filename = filename.replace(char, replacement)
    
    filename = filename.replace(' ', '_')
    
    while '..' in filename:
        filename = filename.replace('..', '.')
    
    if not filename or filename == '.':
        filename = "unnamed_file"
    
    return filename


def apply_pdf_watermark(pdf_path, text="Smart Document Management"):
    """Overlay a transparent background watermark on every page of a PDF."""
    try:
        import fitz
        import math
        from reportlab.pdfgen import canvas
        from reportlab.lib.colors import Color
        
        doc = fitz.open(pdf_path)
        if len(doc) == 0:
            doc.close()
            return
            
        w_docs = {}
        for page_num in range(len(doc)):
            page = doc[page_num]
            rect = page.rect
            width = rect.width
            height = rect.height
            size_key = (width, height)
            
            if size_key not in w_docs:
                temp_w_path = f"{pdf_path}_w_{page_num}.pdf"
                c = canvas.Canvas(temp_w_path, pagesize=(width, height))
                c.saveState()
                c.setFillColor(Color(0.5, 0.5, 0.5, alpha=0.08))
                
                # Calculate diagonal and angle
                diagonal = math.sqrt(width**2 + height**2)
                angle = math.degrees(math.atan2(height, width))
                
                # Calculate dynamic font size to span 85% of diagonal
                font_name = "Helvetica-Bold"
                base_width = c.stringWidth(text, font_name, 10)
                font_size = max(16, int((0.85 * diagonal / base_width) * 10))
                
                c.translate(width / 2.0, height / 2.0)
                c.rotate(angle)
                c.setFont(font_name, font_size)
                c.drawCentredString(0, 0, text)
                c.restoreState()
                c.save()
                
                w_docs[size_key] = (fitz.open(temp_w_path), temp_w_path)
                
            w_doc, temp_w_path = w_docs[size_key]
            # Overlay=False makes it a background watermark (rendered beneath text and drawings)
            page.show_pdf_page(rect, w_doc, 0, overlay=False)
            
        temp_out_path = f"{pdf_path}_watermarked.pdf"
        doc.save(temp_out_path)
        doc.close()
        
        for w_doc, temp_w_path in w_docs.values():
            w_doc.close()
            if os.path.exists(temp_w_path):
                os.remove(temp_w_path)
                
        shutil.move(temp_out_path, pdf_path)
    except Exception as e:
        print(f"Failed to watermark PDF {pdf_path}: {e}")


def apply_docx_watermark(docx_path, text="Smart Document Management"):
    """Inject a transparent background watermark into a Word document's header XML."""
    try:
        import docx
        from docx.oxml import parse_xml
        from docx.oxml.ns import nsdecls
        
        doc = docx.Document(docx_path)
        
        # We append the run containing the shape to the first paragraph of the header.
        # This keeps it in the background of the main document body, but avoids adding
        # a new paragraph (w:p) which introduces layout shifts (like a blank line at the top).
        watermark_xml = (
            f'<w:r {nsdecls("w")}>'
            f'<w:pict>'
            f'<v:shape id="WatermarkObject" type="#_x0000_t136" '
            f'style="position:absolute;left:0;text-align:left;margin-left:0;margin-top:0;'
            f'width:600pt;height:160pt;z-index:-251658240;'
            f'mso-position-horizontal:center;mso-position-horizontal-relative:margin;'
            f'mso-position-vertical:center;mso-position-vertical-relative:margin;rotation:315" '
            f'filled="t" stroked="f" coordsize="21600,21600" xmlns:v="urn:schemas-microsoft-com:vml">'
            f'<v:fill color="#C0C0C0" opacity="0.08"/>'
            f'<v:textpath style="font-family:\'Arial\';font-weight:bold;v-text-kern:t" string="{text}"/>'
            f'</v:shape>'
            f'</w:pict>'
            f'</w:r>'
        )
        
        for section in doc.sections:
            header = section.header
            # Get or create the first paragraph in the header
            if header.paragraphs:
                p = header.paragraphs[0]
            else:
                p = header.add_paragraph()
            
            watermark_el = parse_xml(watermark_xml)
            p._element.append(watermark_el)
            
        doc.save(docx_path)
    except Exception as e:
        print(f"Failed to watermark DOCX {docx_path}: {e}")


def apply_pptx_watermark(pptx_path, text="Smart Document Management"):
    """Add a diagonal semi-transparent watermark text box to every slide of a PPTX."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        prs = Presentation(pptx_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # slide dimensions in inches
        slide_w_in = slide_width / 914400
        slide_h_in = slide_height / 914400
        
        # Watermark textbox size: 90% width, 30% height
        width = slide_width * 0.9
        height = slide_height * 0.3
        left = (slide_width - width) / 2
        top = (slide_height - height) / 2
        
        # Calculate dynamic font size based on slide dimensions
        font_size = Pt(int(min(slide_w_in, slide_h_in) * 7.5))
        
        for slide in prs.slides:
            txBox = slide.shapes.add_textbox(left, top, width, height)
            
            # Send shape to back of the slide visual tree (background)
            try:
                spTree = slide.shapes._spTree
                shape_el = txBox._element
                spTree.remove(shape_el)
                spTree.insert(2, shape_el)
            except:
                pass
                
            tf = txBox.text_frame
            tf.word_wrap = True
            
            # Clear text frame margins for precise centering
            tf.margin_left = Inches(0)
            tf.margin_right = Inches(0)
            tf.margin_top = Inches(0)
            tf.margin_bottom = Inches(0)
            
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = 1  # Center alignment
            
            p.font.size = font_size
            p.font.bold = True
            # Very light gray for transparency simulation (looks transparent on light slides)
            p.font.color.rgb = RGBColor(242, 242, 242)
            p.font.name = 'Arial'
            txBox.rotation = 315
            
        prs.save(pptx_path)
    except Exception as e:
        print(f"Failed to watermark PPTX {pptx_path}: {e}")


def apply_image_watermark(image_path, text="Smart Document Management"):
    """Overlay a rotated semi-transparent watermark on an image using Pillow."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        import math
        
        img = Image.open(image_path)
        orig_format = img.format
        orig_mode = img.mode
        
        img_rgba = img.convert('RGBA')
        txt_layer = Image.new('RGBA', img_rgba.size, (255, 255, 255, 0))
        
        width, height = img.size
        diagonal = math.sqrt(width**2 + height**2)
        
        # Calculate dynamic font size to span 80% of diagonal
        # Arial character width is roughly 0.6 of font size
        font_size = int((0.8 * diagonal) / (len(text) * 0.6))
        font_size = max(16, font_size)
        
        try:
            font = ImageFont.truetype("arial.ttf", font_size)
        except:
            font = ImageFont.load_default()
            
        text_w = len(text) * font_size * 0.6
        text_h = font_size * 1.2
        
        tw = int(text_w) + 60
        th = int(text_h) + 60
        text_img = Image.new('RGBA', (tw, th), (255, 255, 255, 0))
        d = ImageDraw.Draw(text_img)
        
        tx = (tw - text_w) / 2
        ty = (th - text_h) / 2
        # Extremely transparent alpha (20 out of 255, approx 8%)
        d.text((tx, ty), text, font=font, fill=(128, 128, 128, 20))
        
        rotated_txt = text_img.rotate(45, expand=True, resample=Image.BICUBIC)
        
        px = (img.size[0] - rotated_txt.size[0]) // 2
        py = (img.size[1] - rotated_txt.size[1]) // 2
        
        txt_layer.paste(rotated_txt, (px, py), rotated_txt)
        watermarked = Image.alpha_composite(img_rgba, txt_layer)
        
        if orig_format in ['JPEG', 'JPG'] or orig_mode in ['RGB', 'L']:
            watermarked = watermarked.convert(orig_mode if orig_mode != 'RGBA' else 'RGB')
            
        img.close()
        watermarked.save(image_path, format=orig_format)
    except Exception as e:
        print(f"Failed to watermark Image {image_path}: {e}")


def apply_html_watermark(html_path, text="Smart Document Management"):
    """Insert a floating CSS text overlay positioned in the center of the HTML file."""
    try:
        with open(html_path, 'r', encoding='utf-8') as f:
            content = f.read()
            
        watermark_div = f"""
<div style="
    position: fixed;
    top: 50%;
    left: 50%;
    transform: translate(-50%, -50%) rotate(-45deg);
    font-size: calc(6vw + 4vh);
    color: rgba(128, 128, 128, 0.08);
    font-family: sans-serif;
    font-weight: bold;
    pointer-events: none;
    z-index: 9999;
    white-space: nowrap;
    user-select: none;
">{text}</div>
"""
        if '<body>' in content:
            new_content = content.replace('<body>', f'<body>{watermark_div}', 1)
        else:
            new_content = content + watermark_div
            
        with open(html_path, 'w', encoding='utf-8') as f:
            f.write(new_content)
    except Exception as e:
        print(f"Failed to watermark HTML {html_path}: {e}")

